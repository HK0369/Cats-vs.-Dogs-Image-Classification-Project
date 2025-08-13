from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Presentation Setup ---
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# --- Professional Color & Font Scheme ---
colors = {
    "background": RGBColor(255, 255, 255), # White
    "title_text": RGBColor(0, 51, 102),     # Dark Navy Blue
    "body_text": RGBColor(51, 51, 51),      # Dark Gray
    "accent": RGBColor(0, 120, 215)         # Bright Blue
}

def set_font_style(run, size, bold=False, color=colors["body_text"]):
    """Helper function to set font properties."""
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color

# --- Slide Content ---
# This structure contains all the text for the presentation.
slides_data = [
    {
        "layout": "title",
        "title": "Cats vs. Dogs: An End-to-End Image Classification Project",
        "subtitle": "From Data to Deployment\n\nBy N. Harish"
    },
    {
        "layout": "section_header",
        "title": "Phase 1: Data Preparation"
    },
    {
        "layout": "content",
        "title": "Organizing the Raw Data",
        "content": [
            ("The Goal:", "To structure the raw image files into a clean `train/val/test` format that machine learning frameworks can understand."),
            ("The Script:", "`1_data_preparation.py` automates this entire process."),
            ("Key Actions:", [
                "`os.makedirs()`: Creates the necessary directory structure: `data/processed/train`, `data/processed/val`, etc.",
                "`random.shuffle()`: Randomizes the list of training images to ensure the validation set is a fair, unbiased sample.",
                "`shutil.copy()`: Copies the image files from the `raw` source folders to their new destination in the `processed` directory."
            ]),
            ("Reproducibility:", "`random.seed(42)` is used to ensure that the train/validation split is the same every time the script is run, making our results reproducible.")
        ]
    },
    {
        "layout": "section_header",
        "title": "Phase 2: Model Architecture"
    },
    {
        "layout": "content",
        "title": "Building the Convolutional Neural Network (CNN)",
        "content": [
            ("The Goal:", "To define the layers of our neural network, which will learn to extract features from the images."),
            ("The Script:", "`2_initial_training.py` defines the model using Keras's Sequential API."),
            ("The Architecture:", [
                "`Conv2D` Layers: These are the core of the CNN. They act as feature detectors, learning to identify patterns like edges, textures, and shapes. We stack multiple layers (32, 64, 128 filters) to learn increasingly complex features.",
                "`MaxPooling2D` Layers: These downsample the feature maps, reducing computational load and making the detected features more robust to their location in the image.",
                "`Flatten` Layer: Converts the 2D feature map into a 1D vector, preparing it for the classification layers.",
                "`Dense` Layers: These are standard, fully-connected neural network layers that perform the final classification based on the features extracted by the convolutional layers."
            ])
        ]
    },
    {
        "layout": "section_header",
        "title": "Phase 3: Model Training"
    },
    {
        "layout": "content",
        "title": "Teaching the Model to See",
        "content": [
            ("The Goal:", "To train the CNN on our prepared data, allowing it to learn the weights that best differentiate between cats and dogs."),
            ("Data Augmentation (`ImageDataGenerator`):", "To prevent overfitting, we artificially expand our training dataset. The script applies random transformations (rotation, zoom, flips) to the training images, teaching the model to generalize better."),
            # --- FIXED: Corrected the data structure for this item ---
            ("Compilation (`model.compile`):", [
                "Before training, we configure the model with:",
                "A `loss` function (`binary_crossentropy`) to measure prediction error.",
                "An `optimizer` (`Adam`) to update the model's weights.",
                "A `metric` (`accuracy`) to monitor performance."
             ]),
            ("Training (`model.fit`):", "This command starts the training loop, feeding batches of images to the model for a set number of epochs.")
        ]
    },
    {
        "layout": "content",
        "title": "Advanced Training Techniques",
        "content": [
            ("The Goal:", "To improve training efficiency and achieve better results."),
            ("The Script:", "`3_further_training.py` introduces powerful Keras Callbacks."),
            ("Key Callbacks:", [
                "`ModelCheckpoint`: This callback monitors the validation accuracy during training. It automatically saves the model to a file *only* when the accuracy improves. This ensures we always keep the best version of our model.",
                "`EarlyStopping`: This callback also monitors validation accuracy. If the accuracy does not improve for a set number of epochs ('patience'), it automatically stops the training. This saves time and prevents the model from overfitting."
            ]),
            ("Fine-Tuning:", "The script also uses a lower learning rate, a common technique for fine-tuning an already trained model to make smaller, more precise adjustments to its weights.")
        ]
    },
    {
        "layout": "section_header",
        "title": "Phase 4: Frontend Application"
    },
    {
        "layout": "content",
        "title": "Creating the User Interface with Streamlit",
        "content": [
            ("The Goal:", "To create an interactive web application that allows a user to interact with our trained model."),
            ("The Script:", "`app.py` uses the Streamlit library to build the frontend."),
            ("How It Works:", [
                "`st.title()` & `st.markdown()`: Sets up the title and descriptive text on the web page.",
                "`st.file_uploader()`: Creates a widget that allows the user to upload an image file from their computer.",
                "Model Integration: When an image is uploaded, the script opens it with PIL, preprocesses it to match the model's input requirements (resizing to 150x150, normalizing pixels), and feeds it to our loaded `cats_vs_dogs_model.h5`.",
                "Displaying Results: The model's prediction (a probability score) is converted into a human-readable label ('Cat' or 'Dog') and displayed to the user with a confidence percentage."
            ])
        ]
    },
    {
        "layout": "title",
        "title": "Thank You",
        "subtitle": "Questions?"
    }
]

# --- Slide Generation Loop ---
for slide_info in slides_data:
    layout_type = slide_info["layout"]

    if layout_type == "title":
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank layout
        # Add a decorative line at the top
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.5), prs.slide_width, Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["accent"]
        shape.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_title = title_frame.paragraphs[0]
        p_title.alignment = PP_ALIGN.CENTER
        run_title = p_title.add_run()
        run_title.text = slide_info["title"]
        set_font_style(run_title, 54, bold=True, color=colors["title_text"])

        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        p_subtitle = subtitle_frame.paragraphs[0]
        p_subtitle.alignment = PP_ALIGN.CENTER
        run_subtitle = p_subtitle.add_run()
        run_subtitle.text = slide_info["subtitle"]
        set_font_style(run_subtitle, 24, color=colors["body_text"])

    elif layout_type == "section_header":
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank layout
        # Add a background accent shape
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3), prs.slide_width, Inches(3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["title_text"]
        shape.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_title = title_frame.paragraphs[0]
        p_title.alignment = PP_ALIGN.CENTER
        run_title = p_title.add_run()
        run_title.text = slide_info["title"]
        set_font_style(run_title, 48, bold=True, color=RGBColor(255, 255, 255))

    elif layout_type == "content":
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank layout
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1.0))
        p_title = title_shape.text_frame.paragraphs[0]
        run_title = p_title.add_run()
        run_title.text = slide_info["title"]
        set_font_style(run_title, 40, bold=True, color=colors["title_text"])
        
        # Underline for the title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(7), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = colors["accent"]
        line.line.fill.background()

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(7))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for item in slide_info["content"]:
            heading, detail = item
            p_heading = content_frame.add_paragraph()
            p_heading.level = 0
            p_heading.line_spacing = 1.5
            run_heading = p_heading.add_run()
            run_heading.text = heading
            set_font_style(run_heading, 24, bold=True, color=colors["accent"])

            if isinstance(detail, list):
                for sub_item in detail:
                    if isinstance(sub_item, tuple):
                        p_sub_heading = content_frame.add_paragraph()
                        p_sub_heading.level = 1
                        run_sub_heading = p_sub_heading.add_run()
                        run_sub_heading.text = sub_item[0]
                        set_font_style(run_sub_heading, 20, bold=True)
                        for sub_sub_item in sub_item[1]:
                           p_sub_detail = content_frame.add_paragraph()
                           p_sub_detail.level = 2
                           run_sub_detail = p_sub_detail.add_run()
                           run_sub_detail.text = sub_sub_item
                           set_font_style(run_sub_detail, 18)
                    else:
                        p_detail = content_frame.add_paragraph()
                        p_detail.level = 1
                        run_detail = p_detail.add_run()
                        run_detail.text = sub_item
                        set_font_style(run_detail, 20)
            else:
                p_detail = content_frame.add_paragraph()
                p_detail.level = 1
                run_detail = p_detail.add_run()
                run_detail.text = detail
                set_font_style(run_detail, 20)

# --- Save the Presentation ---
pptx_filename = "Cats_vs_Dogs_Project_Detailed.pptx"
prs.save(pptx_filename)

print(f"✅ Detailed PPT saved as '{pptx_filename}'")
